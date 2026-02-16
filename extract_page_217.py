import os
from pptx import Presentation

def extract_slide_217():
    pptx_path = r"e:\AI_Project\MultiAgentDev\inputs\SKB_AICC_콜봇_멀티모달_화면기획서_v1.73_250211.pptx"
    try:
        prs = Presentation(pptx_path)
        slide_idx = 216  # 0-indexed for slide 217
        slide = prs.slides[slide_idx]
        
        print(f"--- Slide {slide_idx + 1} Content ---")
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                # Try multiple encodings for Korean
                text = shape.text
                try:
                    print(text.encode('latin1').decode('cp949'))
                except:
                    print(text)
                    
    except Exception as e:
        print(f"Error: {e}")

if __name__ == "__main__":
    extract_slide_217()
