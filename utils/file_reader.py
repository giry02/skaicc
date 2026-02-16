import os

def read_file_content(file_path):
    """
    Reads the content of a file. Supports text files and PPTX.
    Args:
        file_path (str): The absolute or relative path to the file.
    Returns:
        str: The content of the file, or an error message.
    """
    if not os.path.exists(file_path):
        # Check in 'inputs' directory
        project_root = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
        inputs_path = os.path.join(project_root, "inputs", file_path)
        if os.path.exists(inputs_path):
            file_path = inputs_path
        else:
            return f"[Error] File not found: {file_path}. (Checked absolute path and 'inputs' folder)"
    
    file_ext = os.path.splitext(file_path)[1].lower()

    # PPTX support
    if file_ext == ".pptx":
        try:
            from pptx import Presentation
        except ImportError:
            return "[Error] 'python-pptx' library is not installed. Please run: pip install python-pptx"
        
        try:
            prs = Presentation(file_path)
            text_content = []
            for slide_idx, slide in enumerate(prs.slides):
                slide_text = []
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        slide_text.append(shape.text)
                if slide_text:
                    text_content.append(f"--- Slide {slide_idx + 1} ---\n" + "\n".join(slide_text))
            
            return "\n\n".join(text_content) if text_content else "[Info] No text found in PPTX."
        except Exception as e:
            return f"[Error] Failed to read PPTX file: {e}"

    # Default text reading
    try:
        with open(file_path, "r", encoding="utf-8") as f:
            return f.read()
    except Exception as e:
        return f"[Error] Failed to read file: {e}"
