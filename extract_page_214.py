import sys
from pptx import Presentation

def extract():
    ppt_path = 'inputs/SKB_AICC_콜봇_멀티모달_화면기획서_v1.73_250211.pptx'
    try:
        prs = Presentation(ppt_path)
    except Exception as e:
        print(f"Error loading PPT: {e}")
        return

    # Slide index 213 = Page 214
    slide = prs.slides[213]
    print(f"--- Slide {214} Content ---")
    encodings = ['cp949', 'euc-kr', 'utf-8']
    for shape in slide.shapes:
        if hasattr(shape, "text") and shape.text.strip():
            text = shape.text
            print(f"\nOriginal: {text}")
            for enc in encodings:
                try:
                    # Some shapes might have text that needs to be treated as latin1 before decoding
                    decoded = text.encode('latin1').decode(enc)
                    print(f"[{enc}]: {decoded}")
                except:
                    pass

if __name__ == "__main__":
    extract()
