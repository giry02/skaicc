import os
from pptx import Presentation

def extract_range(start, end):
    pptx_path = r"e:\AI_Project\MultiAgentDev\inputs\SKB_AICC_콜봇_멀티모달_화면기획서_v1.73_250211.pptx"
    try:
        prs = Presentation(pptx_path)
        for i in range(start-1, end):
            slide = prs.slides[i]
            print(f"--- Slide {i + 1} Content ---")
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text = shape.text
                    try:
                        print(text.encode('latin1').decode('cp949'))
                    except:
                        print(text)
    except Exception as e:
        print(f"Error: {e}")

if __name__ == "__main__":
    extract_range(211, 220)
