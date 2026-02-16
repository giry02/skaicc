import sys
from pptx import Presentation

def extract_range(start, end):
    pptx_path = r"e:\AI_Project\MultiAgentDev\inputs\SKB_AICC_콜봇_멀티모달_화면기획서_v1.73_250211.pptx"
    try:
        prs = Presentation(pptx_path)
        for i in range(start-1, end):
            if i >= len(prs.slides):
                break
            slide = prs.slides[i]
            print(f"--- Slide {i + 1} Content ---")
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text = shape.text
                    try:
                        print(text)
                    except UnicodeEncodeError:
                        print(text.encode('ascii', 'ignore').decode('ascii'))
                if shape.has_table:
                    for row in shape.table.rows:
                        for cell in row.cells:
                            print(cell.text)
    except Exception as e:
        print(f"Error: {e}")

def find_keyword(keyword):
    pptx_path = r"e:\AI_Project\MultiAgentDev\inputs\SKB_AICC_콜봇_멀티모달_화면기획서_v1.73_250211.pptx"
    try:
        prs = Presentation(pptx_path)
        for i, slide in enumerate(prs.slides):
            texts = []
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    texts.append(shape.text)
                if shape.has_table:
                    for row in shape.table.rows:
                        for cell in row.cells:
                            texts.append(cell.text)
            full_text = " ".join(texts)
            if keyword.lower() in full_text.lower():
                print(f"--- Keyword Found in Slide {i + 1} ---")
                print(full_text)
    except Exception as e:
        print(f"Error: {e}")

if __name__ == "__main__":
    if len(sys.argv) < 2:
        print("Usage: python extract_slides.py [start] [end] OR python extract_slides.py find [keyword]")
    elif sys.argv[1] == "find":
        find_keyword(sys.argv[2])
    else:
        extract_range(int(sys.argv[1]), int(sys.argv[2]))
